import os
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

from app.analysis.scoring import summarize_for_ppt
from app.viz.charts import bar_scores, bar_components, donut_components, stacked_engagement, sentiment_density, save_fig


def _add_title(prs, title: str, subtitle: str | None = None):
	slide = prs.slides.add_slide(prs.slide_layouts[0])
	slide.shapes.title.text = title
	if subtitle is not None:
		slide.placeholders[1].text = subtitle
	return slide


def _add_image_slide(prs, title: str, image_path: str):
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	slide.shapes.title.text = title
	left = Inches(0.5)
	top = Inches(1.5)
	width = Inches(9)
	slide.shapes.add_picture(image_path, left, top, width=width)
	return slide


def build_ppt(data_dir: str = "data", out_path: str | None = None):
	data_dir = os.path.abspath(data_dir)
	# Load cached CSVs
	posts = pd.read_csv(os.path.join(data_dir, "reddit_posts.csv")) if os.path.exists(os.path.join(data_dir, "reddit_posts.csv")) else pd.DataFrame()
	comments = pd.read_csv(os.path.join(data_dir, "reddit_comments.csv")) if os.path.exists(os.path.join(data_dir, "reddit_comments.csv")) else pd.DataFrame()
	ytv = pd.read_csv(os.path.join(data_dir, "youtube_videos.csv")) if os.path.exists(os.path.join(data_dir, "youtube_videos.csv")) else pd.DataFrame()
	yts = pd.read_csv(os.path.join(data_dir, "youtube_stats.csv")) if os.path.exists(os.path.join(data_dir, "youtube_stats.csv")) else pd.DataFrame()
	trends = pd.read_csv(os.path.join(data_dir, "trends.csv")) if os.path.exists(os.path.join(data_dir, "trends.csv")) else pd.DataFrame()

	# Minimal result reconstruction for charts: expect an existing results CSV or compute quickly
	# For PPT we just craft a light summary using engagement/sentiment/trend columns if present
	# If not present, we skip chart generation
	result_csv = os.path.join(data_dir, "result.csv")
	if os.path.exists(result_csv):
		result = pd.read_csv(result_csv)
	else:
		# Fallback: create a stub result if not available
		result = pd.DataFrame(columns=["place", "score", "engagement", "reddit_engagement", "youtube_engagement", "sentiment", "trend"])

	prs = Presentation()
	date_str = datetime.now().strftime("%d %b %Y")
	_add_title(prs, "X Shoe Store – Chennai Entry Decision", f"Auto-generated on {date_str}")

	# Insight slide
	summary = summarize_for_ppt(result)
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	slide.shapes.title.text = summary["headline"]
	body = slide.placeholders[1].text_frame
	body.clear()
	p = body.paragraphs[0]
	p.text = summary["support"]
	p.font.size = Pt(18)

	# Generate figures and save images
	img_dir = os.path.join(data_dir, "_ppt_images")
	os.makedirs(img_dir, exist_ok=True)

	fig_paths = []
	fig = bar_scores(result)
	if fig:
		path = os.path.join(img_dir, "fit_scores.png")
		save_fig(fig, path)
		fig_paths.append(("Fit Scores: Malls vs High Streets", path))
	fig = bar_components(result)
	if fig:
		path = os.path.join(img_dir, "components_bar.png")
		save_fig(fig, path)
		fig_paths.append(("Signal Components (Bar)", path))
	fig = donut_components(result)
	if fig:
		path = os.path.join(img_dir, "components_donut.png")
		save_fig(fig, path)
		fig_paths.append(("Signal Components (Donut)", path))
	fig = stacked_engagement(result)
	if fig:
		path = os.path.join(img_dir, "engagement_sources.png")
		save_fig(fig, path)
		fig_paths.append(("Engagement Sources by Location", path))
	fig = sentiment_density(result)
	if fig:
		path = os.path.join(img_dir, "sentiment.png")
		save_fig(fig, path)
		fig_paths.append(("Average Sentiment by Location", path))

	# Add image slides
	for title, path in fig_paths:
		_add_image_slide(prs, title, path)

	# Recommendations slide
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	slide.shapes.title.text = "Authentic Launch Recommendations"
	body = slide.placeholders[1].text_frame
	body.clear()
	for b in [
		"Collaborate with local sneaker communities for pop-ups and UGC",
		"Curate influencer runs connecting T Nagar and EA Mall",
		"Limited-edition Chennai colorway drop tied to local sports culture",
		"Pilot store-in-store in high street before flagship scale-up",
	]:
		p = body.add_paragraph()
		p.text = f"• {b}"
		p.level = 0

	# Save
	if out_path is None:
		out_path = os.path.join(data_dir, f"X_Shoe_Chennai_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx")
	prs.save(out_path)
	return out_path


if __name__ == "__main__":
	out = build_ppt()
	print(out)


